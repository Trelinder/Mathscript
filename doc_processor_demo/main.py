import base64
import csv
import io
import logging
import os
from pathlib import Path

import pymupdf
from dotenv import load_dotenv
from fastapi import FastAPI, File, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from fastapi.staticfiles import StaticFiles
from openai import AzureOpenAI
from openpyxl import load_workbook
from pydantic import BaseModel, ConfigDict
from docx import Document
from pptx import Presentation


BASE_DIR = Path(__file__).resolve().parent
load_dotenv(BASE_DIR / ".env")

STATIC_DIR = BASE_DIR / "static"
MAX_UPLOAD_BYTES = 15 * 1024 * 1024
MAX_PDF_PAGES = 8
MAX_TEXT_CHARACTERS = 60_000
SUPPORTED_IMAGE_TYPES = {"image/jpeg", "image/png", "image/webp"}
SUPPORTED_TEXT_EXTENSIONS = {".txt", ".md", ".csv"}
SUPPORTED_OFFICE_EXTENSIONS = {".docx", ".xlsx", ".pptx"}
EXTENSION_MEDIA_TYPES = {
    ".pdf": "application/pdf",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".png": "image/png",
    ".webp": "image/webp",
}

logging.basicConfig(level=os.getenv("LOG_LEVEL", "INFO"))
logger = logging.getLogger("document-processor")


class KeyLineItem(BaseModel):
    model_config = ConfigDict(extra="forbid")

    description: str
    quantity: float | None
    unit_price: float | None
    amount: float | None


class ContractParty(BaseModel):
    model_config = ConfigDict(extra="forbid")

    name: str
    role: str | None


class ContractTerm(BaseModel):
    model_config = ConfigDict(extra="forbid")

    title: str
    summary: str


class DocumentExtraction(BaseModel):
    model_config = ConfigDict(extra="forbid")

    document_category: str
    document_type: str
    entity_name: str | None
    document_number: str | None
    date: str | None
    total_amount: str | None
    key_line_items: list[KeyLineItem]
    effective_date: str | None
    agreement_number: str | None
    parties: list[ContractParty]
    governing_law: str | None
    term_and_renewal: str | None
    payment_terms: str | None
    limitation_of_liability: str | None
    notice_address: str | None
    key_obligations: list[ContractTerm]
    workflow_next_actions: list[str]


app = FastAPI(title="Rapid AI Consultants Document Processing Engine")

configured_origins = [
    origin.strip()
    for origin in os.getenv("CORS_ORIGINS", "*").split(",")
    if origin.strip()
]
app.add_middleware(
    CORSMiddleware,
    allow_origins=configured_origins,
    allow_credentials=False,
    allow_methods=["GET", "POST", "OPTIONS"],
    allow_headers=["*"],
)


def get_openai_client() -> AzureOpenAI:
    endpoint = os.getenv("AZURE_OPENAI_ENDPOINT", "").strip()
    api_key = os.getenv("AZURE_OPENAI_KEY", "").strip()
    api_version = os.getenv("AZURE_OPENAI_API_VERSION", "2024-10-21").strip()

    if not endpoint or not api_key:
        raise HTTPException(
            status_code=503,
            detail="Azure OpenAI credentials are not configured.",
        )

    return AzureOpenAI(
        azure_endpoint=endpoint,
        api_key=api_key,
        api_version=api_version,
        timeout=90.0,
        max_retries=2,
    )


def as_data_url(content: bytes, media_type: str) -> str:
    encoded = base64.b64encode(content).decode("ascii")
    return f"data:{media_type};base64,{encoded}"


def render_pdf_pages(content: bytes) -> list[str]:
    try:
        document = pymupdf.open(stream=content, filetype="pdf")
    except Exception as exc:
        raise HTTPException(status_code=400, detail="The PDF could not be opened.") from exc

    if document.page_count == 0:
        document.close()
        raise HTTPException(status_code=400, detail="The PDF contains no pages.")
    if document.page_count > MAX_PDF_PAGES:
        document.close()
        raise HTTPException(
            status_code=400,
            detail=f"PDFs are limited to {MAX_PDF_PAGES} pages.",
        )

    page_images: list[str] = []
    try:
        for page in document:
            pixmap = page.get_pixmap(matrix=pymupdf.Matrix(1.5, 1.5), alpha=False)
            page_images.append(as_data_url(pixmap.tobytes("jpeg", jpg_quality=85), "image/jpeg"))
    finally:
        document.close()

    return page_images


def truncate_text(text: str) -> str:
    normalized = text.strip()
    if len(normalized) <= MAX_TEXT_CHARACTERS:
        return normalized
    return f"{normalized[:MAX_TEXT_CHARACTERS]}\n\n[Document text truncated for processing]"


def extract_docx_text(content: bytes) -> str:
    try:
        document = Document(io.BytesIO(content))
        sections = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
        for table in document.tables:
            for row in table.rows:
                values = [cell.text.strip() for cell in row.cells]
                if any(values):
                    sections.append(" | ".join(values))
        return truncate_text("\n".join(sections))
    except Exception as exc:
        raise HTTPException(status_code=400, detail="The DOCX file could not be read.") from exc


def extract_xlsx_text(content: bytes) -> str:
    try:
        workbook = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        sections: list[str] = []
        for worksheet in workbook.worksheets:
            sections.append(f"Sheet: {worksheet.title}")
            for row in worksheet.iter_rows(values_only=True):
                values = [str(value).strip() if value is not None else "" for value in row]
                if any(values):
                    sections.append(" | ".join(values))
                if len("\n".join(sections)) >= MAX_TEXT_CHARACTERS:
                    break
        workbook.close()
        return truncate_text("\n".join(sections))
    except Exception as exc:
        raise HTTPException(status_code=400, detail="The XLSX file could not be read.") from exc


def extract_pptx_text(content: bytes) -> str:
    try:
        presentation = Presentation(io.BytesIO(content))
        sections: list[str] = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            sections.append(f"Slide {slide_number}")
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.text.strip():
                    sections.append(shape.text.strip())
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        values = [cell.text.strip() for cell in row.cells]
                        if any(values):
                            sections.append(" | ".join(values))
        return truncate_text("\n".join(sections))
    except Exception as exc:
        raise HTTPException(status_code=400, detail="The PPTX file could not be read.") from exc


def extract_text_document(content: bytes, extension: str) -> str:
    if extension == ".docx":
        text = extract_docx_text(content)
    elif extension == ".xlsx":
        text = extract_xlsx_text(content)
    elif extension == ".pptx":
        text = extract_pptx_text(content)
    elif extension == ".csv":
        try:
            rows = csv.reader(io.StringIO(content.decode("utf-8-sig", errors="replace")))
            text = truncate_text("\n".join(" | ".join(row) for row in rows))
        except csv.Error as exc:
            raise HTTPException(status_code=400, detail="The CSV file could not be read.") from exc
    else:
        text = truncate_text(content.decode("utf-8", errors="replace"))

    if not text:
        raise HTTPException(status_code=422, detail="No readable text was found in this document.")
    return text


def resolve_document_format(file: UploadFile) -> tuple[str, str]:
    extension = Path(file.filename or "").suffix.lower()
    media_type = (file.content_type or "").lower()
    if media_type in SUPPORTED_IMAGE_TYPES:
        return media_type, extension
    if media_type == "application/pdf" or extension == ".pdf":
        return "application/pdf", extension
    if extension in EXTENSION_MEDIA_TYPES:
        return EXTENSION_MEDIA_TYPES[extension], extension
    if extension in SUPPORTED_TEXT_EXTENSIONS | SUPPORTED_OFFICE_EXTENSIONS:
        return "text", extension
    raise HTTPException(
        status_code=415,
        detail="Upload a PDF, image, DOCX, XLSX, PPTX, CSV, TXT, or Markdown document.",
    )


def build_model_content(content: bytes, document_format: str, extension: str) -> list[dict]:
    prompt = (
        "Extract the business document accurately. Set document_category to exactly one of "
        "invoice_order, legal_contract, or general. Preserve identifiers and dates as printed, "
        "and use null when a scalar value is absent. For invoice_order documents, extract "
        "meaningful key_line_items with quantity, unit_price, and amount. For legal_contract "
        "documents, leave key_line_items empty and extract effective_date, agreement_number, "
        "parties, governing_law, term_and_renewal, payment_terms, limitation_of_liability, "
        "notice_address, and key_obligations. Do not invent execution metadata for templates "
        "or boilerplate agreements. Recommend concise downstream workflow actions appropriate "
        "to the document."
    )
    if document_format == "text":
        extracted_text = extract_text_document(content, extension)
        return [{"type": "text", "text": f"{prompt}\n\nDocument content:\n{extracted_text}"}]

    image_urls = (
        render_pdf_pages(content)
        if document_format == "application/pdf"
        else [as_data_url(content, document_format)]
    )
    vision_content: list[dict] = [{"type": "text", "text": prompt}]
    vision_content.extend(
        {
            "type": "image_url",
            "image_url": {"url": image_url, "detail": "high"},
        }
        for image_url in image_urls
    )
    return vision_content


@app.get("/health")
async def health() -> dict[str, str]:
    return {"status": "ok", "model": os.getenv("AZURE_OPENAI_DEPLOYMENT", "gpt-5.4-mini")}


@app.post("/api/process-document", response_model=DocumentExtraction)
async def process_document(file: UploadFile = File(...)) -> DocumentExtraction:
    document_format, extension = resolve_document_format(file)

    content = await file.read(MAX_UPLOAD_BYTES + 1)
    await file.close()
    if not content:
        raise HTTPException(status_code=400, detail="The uploaded file is empty.")
    if len(content) > MAX_UPLOAD_BYTES:
        raise HTTPException(status_code=413, detail="Files are limited to 15 MB.")

    try:
        completion = get_openai_client().chat.completions.create(
            model=os.getenv("AZURE_OPENAI_DEPLOYMENT", "gpt-5.4-mini"),
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You are a precise document-processing engine. Return only data "
                        "that conforms to the supplied JSON schema. Never invent values."
                    ),
                },
                {
                    "role": "user",
                    "content": build_model_content(content, document_format, extension),
                },
            ],
            response_format={
                "type": "json_schema",
                "json_schema": {
                    "name": "document_extraction",
                    "strict": True,
                    "schema": DocumentExtraction.model_json_schema(),
                },
            },
        )
        message = completion.choices[0].message
        if message.refusal:
            raise HTTPException(status_code=422, detail=f"Document declined: {message.refusal}")
        if not message.content:
            raise ValueError("Azure OpenAI returned an empty response.")
        return DocumentExtraction.model_validate_json(message.content)
    except HTTPException:
        raise
    except Exception as exc:
        logger.exception("Document processing failed")
        raise HTTPException(
            status_code=502,
            detail="The document could not be processed by the AI service.",
        ) from exc


@app.get("/", include_in_schema=False)
async def dashboard() -> FileResponse:
    return FileResponse(STATIC_DIR / "index.html")


app.mount("/static", StaticFiles(directory=STATIC_DIR), name="static")